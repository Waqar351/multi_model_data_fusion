from __future__ import annotations

import csv
import sys
from pathlib import Path

sys.path.insert(0, str(Path('.presentation_deps').resolve()))

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
RESULTS = ROOT / 'results'
OUT = ROOT / 'explainability_comparison_presentation.pptx'
CSV_OUT = ROOT / 'explainability_comparison_results.csv'

MODELS = [
    ('M7 Late Fusion', 'M7_Late_Fusion_Complete_Explainability'),
    ('M11 Scalar Gate', 'M11_Gate_Scalar_Complete_Explainability'),
    ('M11 Vector Gate', 'M11_Gate_Vector_only_Explainability'),
    ('M33 Hierarchical', 'M33_Hierarchical_Fusion_Complete_Explainability'),
    ('Original GMU', 'Original_GMU_Only_Complete_Explainability'),
    ('GMU–GNN', 'GMU_GNN_Complete_Explainability'),
]

METRICS = {
    'M7 Late Fusion': (0.00045953551307320595, 0.021436779447323844, 0.012592143379151821, 0.8416081070899963),
    'M11 Scalar Gate': (5.324837184161879e-05, 0.007297148199236383, 0.0037909741513431072, 0.9816464185714722),
    'M11 Vector Gate': (5.5545708164572716e-05, 0.007452899312654956, 0.00409318320453167, 0.9808546304702759),
    'M33 Hierarchical': (0.00012352352496236563, 0.011114113773142942, 0.0061009968630969524, 0.9574241042137146),
    'Original GMU': (5.7567078329157084e-05, 0.0075872971688973065, 0.0036113597452640533, 0.9801578521728516),
    'GMU–GNN': (7.590175664518028e-05, 0.008712161422125986, 0.004443921148777008, 0.9738383293151855),
}

COLORS = ['315C6B', '2F80ED', '8B5CF6', 'E08A1E', '2AA876', 'D1495B']
BG = RGBColor(247, 248, 250)
INK = RGBColor(28, 35, 48)
MUTED = RGBColor(91, 101, 116)
WHITE = RGBColor(255, 255, 255)


def rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font='Aptos', valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def base_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BG
    add_text(slide, title, .55, .28, 12.2, .45, 25, INK, True)
    if subtitle:
        add_text(slide, subtitle, .57, .77, 12.0, .34, 10.5, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.55), Inches(1.08), Inches(12.2), Inches(.025))
    line.fill.solid(); line.fill.fore_color.rgb = rgb('D9DEE7'); line.line.fill.background()
    return slide


def footer(slide, number):
    add_text(slide, 'Source: saved outputs in results/<notebook name> • chronological test split • seed 42',
             .58, 7.18, 11.7, .18, 7.5, MUTED)
    add_text(slide, str(number), 12.35, 7.15, .35, .2, 8, MUTED, align=PP_ALIGN.RIGHT)


def fit_image(slide, path: Path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    slide.shapes.add_picture(str(path), Inches(x + (w-dw)/2), Inches(y + (h-dh)/2),
                             width=Inches(dw), height=Inches(dh))


def comparison_grid(prs, title, subtitle, files, takeaway):
    slide = base_slide(prs, title, subtitle)
    card_w, card_h = 4.03, 2.72
    for i, ((label, folder), filename) in enumerate(zip(MODELS, files)):
        x = .55 + (i % 3) * 4.12; y = 1.22 + (i // 3) * 2.83
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = rgb('DCE1E8')
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(.07), Inches(card_h))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = rgb(COLORS[i]); stripe.line.fill.background()
        add_text(slide, label, x+.18, y+.08, card_w-.28, .24, 11, rgb(COLORS[i]), True)
        if filename:
            p = RESULTS / folder / filename
            if p.exists():
                fit_image(slide, p, x+.12, y+.40, card_w-.24, card_h-.5)
            else:
                add_text(slide, 'Export not found', x+.3, y+1.15, card_w-.6, .35, 13, MUTED, align=PP_ALIGN.CENTER)
        else:
            add_text(slide, 'Not applicable\n(no adaptive gate for this architecture)', x+.35, y+1.03, card_w-.7, .65,
                     13, MUTED, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, takeaway, .72, 6.91, 11.9, .24, 9.5, INK, True, PP_ALIGN.CENTER)
    footer(slide, len(prs.slides))
    return slide


def make_deck():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    prs.core_properties.title = 'Cross-model explainability comparison'
    prs.core_properties.subject = 'Comparison of six fusion notebooks and their saved outputs'

    # 1 Title
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb('17202D')
    add_text(s, 'Fusion-model explainability', .72, 1.05, 11.9, .7, 36, WHITE, True)
    add_text(s, 'A cross-notebook comparison of performance, modality use, interventions, graph reliance and representations',
             .75, 1.87, 11.4, 1.0, 20, rgb('DDE6F1'))
    for i, (label, _) in enumerate(MODELS):
        x = .78 + (i % 3)*4.05; y = 3.25 + (i//3)*1.02
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.68), Inches(.72))
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb(COLORS[i]); sh.line.fill.background()
        add_text(s, label, x+.15, y+.18, 3.38, .3, 15, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, 'Prepared from saved notebook outputs • 17 August 2026', .78, 6.75, 11.5, .3, 10, rgb('AAB7C7'))

    # 2 scope
    s=base_slide(prs, 'Comparison design', 'Equivalent analyses are aligned; architecture-specific analyses are marked not applicable.')
    cards=[
        ('Common evaluation', 'Test metrics, training diagnostics, prediction diagnostics, ablation/intervention, permutation importance, uncertainty and embeddings.'),
        ('Adaptive fusion', 'Gate health, feature-wise specialization, gate–crime/error association, and effective gated activation apply to scalar/vector/GMU models.'),
        ('Graph-specific evidence', 'Graph interventions and neighborhood context apply only to graph models; Original GMU is the explicit no-graph comparator.'),
        ('Interpretive rule', 'Raw gate coefficients are not treated as contribution. Effective activation or occlusion is used when available.'),
    ]
    for i,(h,b) in enumerate(cards):
        x=.75+(i%2)*6.15; y=1.42+(i//2)*2.28
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.65), Inches(1.78))
        sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=rgb('DCE1E8')
        add_text(s,h,x+.25,y+.22,5.15,.32,17,rgb(COLORS[i]),True)
        add_text(s,b,x+.25,y+.65,5.12,.85,13,MUTED)
    add_text(s,'All numerical claims in this deck come from the six current results folders, not results_old.',.8,6.45,11.7,.35,12,INK,True,PP_ALIGN.CENTER)
    footer(s,len(prs.slides))

    # 3 architecture
    s=base_slide(prs,'Architecture map','Fusion location and graph usage explain why some analyses are not directly interchangeable.')
    arch=[
        ('Late Fusion','Static GNN + Dynamic GNN → late fusion','Graph • no explicit gate'),
        ('M11 Scalar','Global self/neighbor scalar mixing','Graph • 2 fixed learned coefficients'),
        ('M11 Vector','Adaptive feature-wise self/neighbor gates','Graph • case-dependent vector gates'),
        ('Hierarchical','Separate graph encoders → fusion → joint graph','Graph • no explicit modality gate'),
        ('Original GMU','Single adaptive GMU → regression','No graph • feature-wise gate'),
        ('GMU–GNN','Single adaptive GMU → graph model','Graph • feature-wise gate'),
    ]
    for i,(h,b,c) in enumerate(arch):
        x=.58+(i%3)*4.16; y=1.28+(i//3)*2.68
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.88), Inches(2.18))
        sh.fill.solid(); sh.fill.fore_color.rgb=WHITE; sh.line.color.rgb=rgb(COLORS[i])
        add_text(s,h,x+.22,y+.18,3.45,.35,17,rgb(COLORS[i]),True)
        add_text(s,b,x+.22,y+.68,3.42,.72,14,INK)
        add_text(s,c,x+.22,y+1.60,3.42,.3,11,MUTED,True)
    footer(s,len(prs.slides))

    # 4 performance table
    s=base_slide(prs,'Test performance','Lower is better for MSE/RMSE/MAE; higher is better for R².')
    rows=sorted(METRICS.items(), key=lambda kv: kv[1][0])
    x0=.75; widths=[3.2,2.1,2.1,2.1,1.9]; y0=1.35; rh=.68
    headers=['Model','MSE','RMSE','MAE','R²']
    for j,(h,w) in enumerate(zip(headers,widths)):
        x=x0+sum(widths[:j]); sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y0), Inches(w), Inches(rh))
        sh.fill.solid(); sh.fill.fore_color.rgb=rgb('26384A'); sh.line.color.rgb=WHITE
        add_text(s,h,x+.08,y0+.19,w-.16,.25,12,WHITE,True,PP_ALIGN.CENTER)
    for i,(name,vals) in enumerate(rows):
        y=y0+(i+1)*rh; color=COLORS[[m[0] for m in MODELS].index(name)]
        texts=[name,f'{vals[0]:.3e}',f'{vals[1]:.5f}',f'{vals[2]:.5f}',f'{vals[3]:.4f}']
        for j,(t,w) in enumerate(zip(texts,widths)):
            x=x0+sum(widths[:j]); sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(rh))
            sh.fill.solid(); sh.fill.fore_color.rgb=WHITE if i%2==0 else rgb('F1F4F7'); sh.line.color.rgb=rgb('DCE1E8')
            add_text(s,t,x+.08,y+.19,w-.16,.25,12,rgb(color) if j==0 else INK,j==0,PP_ALIGN.CENTER)
    add_text(s,'M11 Scalar leads MSE/RMSE/R²; Original GMU leads MAE. The ranking gap between the top three is small relative to M7.',
             .8,6.15,11.7,.58,13,INK,True,PP_ALIGN.CENTER)
    footer(s,len(prs.slides))

    comparison_grid(prs,'Training diagnostics','Loss curves and gate/coefficient histories from each experiment.',[
        'train_validation_loss_M7_Late_Fusion_Complete_Explainability.png','scalar_training_validation_coefficients.png',
        'm11_vector_gate_training_history.png','hierarchical_training_diagnostics.png','original_gmu_chronological_training.png',
        'gmu_gnn_chronological_training.png'], 'Use these curves to assess convergence; final test metrics remain the fair ranking criterion.')
    comparison_grid(prs,'Prediction diagnostics','Observed-vs-predicted and residual behavior on the test period.',[
        None,'scalar_prediction_diagnostics.png','m11_prediction_diagnostics.png','hierarchical_prediction_diagnostics.png',
        'gmu_prediction_diagnostics.png','gmu_prediction_diagnostics.png'], 'The strongest models cluster most tightly around the identity line; M7 has no matching exported diagnostic.')
    comparison_grid(prs,'Modality / branch contribution','Effective activation or branch occlusion; these are more faithful than raw gate means.',[
        'late_fusion_effective_contribution.png','scalar_health_effective_roles.png','professional_effective_activation_contribution.png',
        'hierarchical_global_contribution.png','gmu_effective_activation_contribution.png','gmu_effective_activation_contribution.png'],
        'Dynamic evidence is indispensable in every model, while effective static share ranges from ~6% (M7) to ~94% (scalar M11).')
    comparison_grid(prs,'Gate health','Distribution, entropy and saturation of learned gates or coefficients.',[
        None,'scalar_health_effective_roles.png','professional_gate_health.png',None,'gmu_professional_gate_health.png','gmu_professional_gate_health.png'],
        'Vector/GMU gates are unsaturated overall; scalar coefficients are near 0.5 but their effective activations are strongly static-heavy.')
    comparison_grid(prs,'Feature / channel specialization','Which latent dimensions preferentially carry static or dynamic information.',[
        'late_fusion_channel_specialization.png',None,'professional_featurewise_gate_specialization.png','hierarchical_encoder_channels.png',
        'gmu_featurewise_gate_specialization.png','gmu_featurewise_gate_specialization.png'],
        'Specialization is architecture-specific: channels for branch models, feature-wise gates for vector/GMU models.')
    comparison_grid(prs,'Association with crime level','Fusion behavior versus input and future crime.',[
        'late_fusion_crime_associations.png','scalar_effective_activation_crime.png','professional_gate_crime_associations.png',
        'hierarchical_share_crime.png','gmu_gate_crime_associations.png','gmu_gate_crime_associations.png'],
        'Adaptive gates change systematically with crime; M11 vector separates self and neighbor roles in opposite directions.')
    comparison_grid(prs,'Spatial and temporal stability','Node-level heterogeneity and time-level stability of fusion behavior.',[
        'late_fusion_spatial_temporal_stability.png','scalar_effective_spatial_temporal.png','professional_spatial_temporal_gate_stability.png',
        'hierarchical_spatial_temporal.png','gmu_spatial_temporal_gate_stability.png','gmu_spatial_temporal_gate_stability.png'],
        'Mean behavior is temporally stable, but node-level variation is meaningful—especially for hierarchical and vector gating.')
    comparison_grid(prs,'Fusion behavior versus error','Whether gate/share changes are associated with larger absolute prediction errors.',[
        'late_fusion_share_error.png','scalar_effective_balance_error.png','professional_gate_error_relationship.png','hierarchical_share_error.png',
        'gmu_gate_error_relationship.png','gmu_gate_error_relationship.png'],
        'Associations are descriptive, not causal; intervention slides provide the stronger faithfulness evidence.')
    comparison_grid(prs,'Intervention and ablation faithfulness','Performance after forcing, removing, shuffling or reassigning modalities.',[
        'late_fusion_branch_interventions.png','scalar_interventions_corrected.png','professional_gate_intervention_faithfulness.png',
        'hierarchical_branch_interventions.png','gmu_gate_intervention_faithfulness.png','gmu_gate_intervention_faithfulness.png'],
        'Learned fusion is strongly preferred; dynamic removal is catastrophic, and vector-gate reassignment is exceptionally damaging.')
    comparison_grid(prs,'Input permutation importance','Model-agnostic feature importance measured as ΔMSE after permutation.',[
        'late_fusion_input_permutation.png','scalar_input_permutation.png','professional_input_permutation_importance.png',
        'hierarchical_input_permutation.png','gmu_input_permutation_importance.png','gmu_input_permutation_importance.png'],
        'Recent crime lags dominate across models; static features add smaller but model-dependent increments.')
    comparison_grid(prs,'Graph reliance','Counterfactual edge removal or hierarchy-stage bypass.',[
        'late_fusion_graph_intervention.png','scalar_exact_hop_context.png',None,'hierarchy_level_interventions.png',None,
        'gmu_gnn_graph_intervention.png'],
        'Edges help M11 scalar, M33 and GMU–GNN; M7 improves by 79.7% MSE when edges are removed, signaling harmful graph use.')
    comparison_grid(prs,'Bootstrap uncertainty','Time-cluster bootstrap confidence intervals for key quantities.',[
        'late_fusion_bootstrap_ci.png',None,'professional_gate_bootstrap_ci.png','hierarchical_bootstrap_ci.png',
        'gmu_gate_bootstrap_ci.png','gmu_gate_bootstrap_ci.png'],
        'Intervals are narrow because test cases are numerous; time-cluster resampling is preferable to independent case resampling.')
    comparison_grid(prs,'Embedding projections','t-SNE views of learned representations; M33 shows both hierarchy stages.',[
        'late_fusion_embedding_tsne_robust.png','scalar_embedding_tsne.png','m11_embedding_tsne.png','hierarchical_stage_tsne.png',
        'gmu_embedding_tsne.png','gmu_embedding_tsne.png'],
        'Embedding geometry is qualitative evidence only; it should support—not replace—metric and intervention results.')
    comparison_grid(prs,'High-crime node / neighborhood context','Representative local behavior and exact-hop spatial context.',[
        'M7_Late_Fusion_Complete_Explainability_high_crime_node_397.png','scalar_exact_hop_context.png',
        'm11_high_crime_node_neighborhood.png','hierarchical_high_crime_neighborhood.png','gmu_high_crime_node_neighborhood.png',
        'gmu_high_crime_node_neighborhood.png'], 'Original GMU has no message passing; its neighborhood plot is contextual rather than a graph-reliance explanation.')

    # Synthesis
    s=base_slide(prs,'Overall synthesis','What the aligned evidence supports—and what it does not.')
    points=[
        ('Best predictive fit','M11 Scalar: MSE 5.325×10⁻⁵ and R² 0.9816. Original GMU has the lowest MAE (0.00361).'),
        ('Most expressive gate evidence','M11 Vector: feature-wise self/neighbor specialization and opposite crime associations, but slightly weaker prediction than scalar M11.'),
        ('Graph value is model-dependent','Edges materially help scalar M11, M33 and GMU–GNN; they hurt M7 in the saved edge-removal counterfactual.'),
        ('Dynamic signal is essential','Dynamic-only is consistently far stronger than static-only; removing/shuffling dynamic branches produces the largest failures.'),
        ('Raw gates can mislead','Scalar coefficients near 0.5 coexist with ~94% effective static activation; interpret contribution using activations/occlusion.'),
        ('Recommended reporting','Lead with performance + intervention + permutation evidence; use gate plots, stability and t-SNE as supporting interpretation.'),
    ]
    for i,(h,b) in enumerate(points):
        y=1.30+i*.86
        circ=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(.78), Inches(y+.02), Inches(.34), Inches(.34))
        circ.fill.solid(); circ.fill.fore_color.rgb=rgb(COLORS[i]); circ.line.fill.background()
        add_text(s,str(i+1),.78,y+.075,.34,.18,9,WHITE,True,PP_ALIGN.CENTER)
        add_text(s,h,1.30,y,2.55,.35,14,rgb(COLORS[i]),True)
        add_text(s,b,3.72,y,8.75,.52,12.5,INK)
    footer(s,len(prs.slides))

    prs.save(OUT)
    with CSV_OUT.open('w', newline='', encoding='utf-8-sig') as f:
        w=csv.writer(f); w.writerow(['model','MSE','RMSE','MAE','R2','MSE_rank'])
        ranks={name:i+1 for i,(name,_) in enumerate(rows)}
        for name,_ in MODELS: w.writerow([name,*METRICS[name],ranks[name]])
    print(f'Wrote {OUT} ({len(prs.slides)} slides)')
    print(f'Wrote {CSV_OUT}')


if __name__ == '__main__':
    make_deck()
